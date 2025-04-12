import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

from . import PPTX_STYLE


def hex_to_rgb(hex_color):
    if hex_color.startswith("#") and len(hex_color) == 7:
        hex_color = hex_color.lstrip("#")
        return RGBColor(int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16))
    else:
        print(f"[!] Некорректный цвет: {hex_color}, fallback to red")
        return RGBColor(255, 0, 0)


def estimate_lines(text: str, wrap_limit=80):
    lines = 0
    for paragraph in text.split("\n"):
        lines += max(1, len(paragraph) // wrap_limit + 1)
    return lines


def get_titul(json_data):
    return {
        "text": json_data["titul"]["text"],
        "image": json_data["titul"].get("image_prompt", "")
    }


def get_main(json_data):
    main = json_data["main"]
    return {
        "plan": main.get("plan", []),
        "color_header": main["style"]["color_header"],
        "color_text": main["style"]["color_text"],
        "image": main.get("image_prompt", "")
    }


def get_body(json_data):
    slides = []
    for slide in json_data["body"]:
        title = slide["slide"]["title"]
        content_blocks = slide["slide"]["content"]
        image = slide.get("image", "")
        slides.append((title, content_blocks, image))
    return slides


class GeneratePresentation:
    def __init__(self, titul, main, body):
        self.titul = titul
        self.main = main
        self.body = body
        self.pr = Presentation()
        self.pr.slide_width = Inches(PPTX_STYLE["slide_width"])
        self.pr.slide_height = Inches(PPTX_STYLE["slide_height"])

    def add_titul(self):
        slide = self.pr.slides.add_slide(self.pr.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
        tf = title_box.text_frame
        tf.text = self.titul["text"]
        p = tf.paragraphs[0]
        style = PPTX_STYLE["title_font"]
        p.font.size = Pt(style["size"])
        p.font.bold = style["bold"]
        p.alignment = getattr(PP_ALIGN, style["align"].upper())
        return slide

    def add_plan(self):
        slide = self.pr.slides.add_slide(self.pr.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
        tf_title = title_box.text_frame
        tf_title.text = "План"
        p_title = tf_title.paragraphs[0]
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.LEFT

        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for item in self.main["plan"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(24)
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.LEFT

    def get_image_path(self, image_path):
        return image_path if os.path.exists(image_path) else None

    def add_body(self):
        for title, content_blocks, image in self.body:
            if not content_blocks:
                continue

            slide = self.pr.slides.add_slide(self.pr.slide_layouts[6])

            # Title
            s = PPTX_STYLE["slide_title"]
            title_box = slide.shapes.add_textbox(Inches(s["left"]), Inches(s["top"]), Inches(s["width"]), Inches(s["height"]))
            tf_title = title_box.text_frame
            tf_title.text = title
            p_title = tf_title.paragraphs[0]
            p_title.font.size = Pt(s["font_size"])
            p_title.font.bold = s["bold"]
            p_title.alignment = getattr(PP_ALIGN, s["align"].upper())

            total_lines = 0
            for block in content_blocks:
                text = block.get("text", "")
                total_lines += 1  # subtitle
                total_lines += estimate_lines(text, wrap_limit=80)

            line_height_pt = PPTX_STYLE["text"]["font_size"] + 8
            total_height_pt = total_lines * line_height_pt
            available_height = PPTX_STYLE["content_box"]["height"] * 72
            top_offset_pt = max((available_height - total_height_pt) // 2, 0)
            top_offset_inches = top_offset_pt / 72

            c = PPTX_STYLE["content_box"]
            img = PPTX_STYLE["image"]
            text_width = c["width"] * 0.55
            content_box = slide.shapes.add_textbox(
                Inches(c["left"]),
                Inches(c["top"] + top_offset_inches),
                Inches(text_width),
                Inches(c["height"])
            )
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            if image:
                for block in content_blocks:
                    sub = PPTX_STYLE["subtitle"]
                    p_sub = tf_content.add_paragraph()
                    p_sub.text = block.get("subtitle", "")
                    p_sub.font.size = Pt(sub["font_size"])
                    p_sub.font.bold = sub["bold"]
                    p_sub.font.name = sub["font_name"]
                    p_sub.font.color.rgb = hex_to_rgb(self.main["color_text"])
                    p_sub.alignment = getattr(PP_ALIGN, sub["align"].upper())
                    p_sub.space_after = Pt(4)

                    txt = PPTX_STYLE["text"]
                    p_txt = tf_content.add_paragraph()
                    p_txt.text = block.get("text", "")
                    p_txt.font.size = Pt(txt["font_size"])
                    p_txt.font.name = txt["font_name"]
                    p_txt.font.color.rgb = hex_to_rgb(self.main["color_text"])
                    p_txt.alignment = getattr(PP_ALIGN, txt["align"].upper())
                    p_txt.space_after = Pt(18)


                    image_path = self.get_image_path(image)
                    if image_path:
                        try:
                            text_width_inches = text_width
                            text_gap = Inches(0.5)

                            img_width = Inches(img["width"] * 0.42)
                            img_height = img_width

                            slide_height = self.pr.slide_height
                            img_top = (slide_height - img_height) / 2

                            img_left = Inches(c["left"]) + Inches(text_width_inches) + text_gap + Inches(0.3)

                            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
                        except Exception as e:
                            print(f"[!] Ошибка вставки изображения: {e}")

            else:
                for block in content_blocks:
                    sub = PPTX_STYLE["subtitle"]
                    p_sub = tf_content.add_paragraph()
                    p_sub.text = block.get("subtitle", "")
                    p_sub.font.size = Pt(sub["font_size"])
                    p_sub.font.bold = sub["bold"]
                    p_sub.font.name = sub["font_name"]
                    p_sub.font.color.rgb = hex_to_rgb(self.main["color_text"])
                    p_sub.alignment = getattr(PP_ALIGN, sub["align"].upper())
                    p_sub.space_after = Pt(4)

                    txt = PPTX_STYLE["text"]
                    p_txt = tf_content.add_paragraph()
                    p_txt.text = block.get("text", "")
                    p_txt.font.size = Pt(txt["font_size"])
                    p_txt.font.name = txt["font_name"]
                    p_txt.font.color.rgb = hex_to_rgb(self.main["color_text"])
                    p_txt.alignment = getattr(PP_ALIGN, txt["align"].upper())
                    p_txt.space_after = Pt(18)


    def save(self, filename):
        self.pr.save(filename)
        return filename


def generate_ppt_from_json(data: dict, filename: str):
    prs = GeneratePresentation(get_titul(data), get_main(data), get_body(data))
    prs.add_titul()
    prs.add_plan()
    prs.add_body()
    prs.save(filename)
    return filename


if __name__ == "__main__":
    with open("hi.json", "r", encoding="utf-8") as f:
        data = json.load(f)
    generate_ppt_from_json(data, "output_presentation.pptx")
